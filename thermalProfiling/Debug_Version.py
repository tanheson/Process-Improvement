# --------------------------------------------------------------
#  DEBUG VERSION – will tell you why nothing is shown in PPTX
# --------------------------------------------------------------
import os
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

# ------------------- CONFIG -------------------
root_dir = r'U:\NVL\HX\A0\results_experimental\ThermalProfile'
# ------------------------------------------------

subfolder_data = {}

# ---------- 1. FILE PROCESSING ----------
def process_file(file_path, subfolder_name):
    print(f"\n>>> Trying to read: {file_path}")
    try:
        if os.path.basename(file_path).startswith('~$'):
            print("    → Skipped (temporary Excel file)")
            return

        df = pd.read_excel(file_path, sheet_name='SearchVoltage Results', engine='openpyxl')
        required_cols = [
            'run_order',
            'cmvprofiling.Max_DTS_Profile_max',
            'cmvprofiling.Intec_TC_Profile_max',
            'cmvprofiling.Intec_FB_Profile_max'
        ]

        missing = [c for c in required_cols if c not in df.columns]
        if missing:
            print(f"    → Missing columns {missing} – SKIPPED")
            return

        valid_data = df[required_cols].dropna()
        if valid_data.empty:
            print("    → No rows with all required columns – SKIPPED")
            return

        # Initialise container
        if subfolder_name not in subfolder_data:
            subfolder_data[subfolder_name] = {'run_x': [], 'dts_y': [], 'tc_y': [], 'fb_y': []}

        run_x = valid_data['run_order'].astype(int).tolist()
        dts_y = pd.to_numeric(valid_data['cmvprofiling.Max_DTS_Profile_max'], errors='coerce').dropna().astype(int).tolist()
        tc_y  = pd.to_numeric(valid_data['cmvprofiling.Intec_TC_Profile_max'],  errors='coerce').dropna().astype(int).tolist()
        fb_y  = pd.to_numeric(valid_data['cmvprofiling.Intec_FB_Profile_max'],  errors='coerce').dropna().astype(int).tolist()

        min_len = min(len(run_x), len(dts_y), len(tc_y), len(fb_y))
        if min_len == 0:
            print("    → One of the series is empty after cleaning – SKIPPED")
            return

        subfolder_data[subfolder_name]['run_x'].extend(run_x[:min_len])
        subfolder_data[subfolder_name]['dts_y'].extend(dts_y[:min_len])
        subfolder_data[subfolder_name]['tc_y'].extend(tc_y[:min_len])
        subfolder_data[subfolder_name]['fb_y'].extend(fb_y[:min_len])

        print(f"    → Added {min_len} rows → total now {len(subfolder_data[subfolder_name]['run_x'])}")

    except Exception as e:
        print(f"    → EXCEPTION: {e}")

# ---------- 2. STATISTICS ----------
def calculate_stats(series, ranges):
    stats = {f"{high}-{low}": 0 for high, low in ranges}
    for temp in series:
        for high, low in ranges:
            if low <= temp <= high:
                stats[f"{high}-{low}"] += 1
                break
    return stats

# ---------- 3. LATEST HotVmin TIMESTAMP ----------
def get_latest_hotvmin_path(hotvmin_root):
    """
    hotvmin_root = .../D3/HotVmin   (or D4)
    Returns full path to the newest timestamp folder or None.
    """
    if not os.path.isdir(hotvmin_root):
        print(f"    HotVmin folder not found: {hotvmin_root}")
        return None

    candidates = []
    for name in os.listdir(hotvmin_root):
        full = os.path.join(hotvmin_root, name)
        if not os.path.isdir(full):
            continue
        try:
            ts = datetime.strptime(name, "%Y.%m.%d_%H.%M.%S")
            candidates.append((name, ts, full))
        except ValueError:
            continue

    if not candidates:
        print(f"    No timestamped sub-folders inside {hotvmin_root}")
        return None

    latest_name, latest_ts, latest_path = max(candidates, key=lambda x: x[1])
    print(f"    Latest HotVmin timestamp → {latest_name} ({latest_ts})")
    return latest_path

# -------------------- MAIN --------------------
try:
    main_folders = [f for f in os.listdir(root_dir) if 'D3' in f or 'D4' in f]
    if not main_folders:
        raise RuntimeError("No D3/D4 folders found under root_dir")

    for main in main_folders:
        main_path = os.path.join(root_dir, main)
        print(f"\n=== PROCESSING MAIN FOLDER: {main} ===")
        subfolder_data.clear()

        # ---- list immediate sub-folders ----
        subfolders = [d for d in os.listdir(main_path) if os.path.isdir(os.path.join(main_path, d))]
        print(f"  Sub-folders found: {subfolders}")

        for subfolder in subfolders:
            subfolder_path = os.path.join(main_path, subfolder)
            print(f"\n  → Handling sub-folder: {subfolder}")

            if subfolder == "HotVmin":
                latest_path = get_latest_hotvmin_path(subfolder_path)
                if latest_path:
                    subfolder_name = os.path.basename(latest_path)   # e.g. 2025.07.30_21.48.51
                    print(f"    Walking only: {latest_path}")
                    for dirpath, _, filenames in os.walk(latest_path):
                        for fn in filenames:
                            if fn.endswith('.xlsx') and ('HotVmin' in fn or 'HotGNG' in fn):
                                process_file(os.path.join(dirpath, fn), subfolder_name)
                else:
                    print("    No HotVmin data processed")
            else:
                # ---- normal sub-folders (HotGNG, etc.) ----
                for dirpath, _, filenames in os.walk(subfolder_path):
                    for fn in filenames:
                        if fn.endswith('.xlsx') and ('HotVmin' in fn or 'HotGNG' in fn):
                            process_file(os.path.join(dirpath, fn), subfolder)

        # ---- AFTER ALL FILES ARE READ ----
        print(f"\n  Collected sub-folders with data: {list(subfolder_data.keys())}")
        if not subfolder_data:
            print("  No data at all → PPTX will be empty")
            continue

        # ---- PPTX CREATION ----
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # blank

        temp_ranges = [(130,120),(119,110),(109,105),(104,100),(99,95),
                       (94,90),(89,85),(84,80),(70,60)]
        temp_ranges.sort(key=lambda x: x[0], reverse=True)

        for subfolder, data in subfolder_data.items():
            if not any(data.values()):
                print(f"  Skipping empty slide for {subfolder}")
                continue

            # ---- plot ----
            plt.figure(figsize=(18,8))
            plt.scatter(data['run_x'], data['dts_y'], c='blue', label='Max_DTS')
            plt.scatter(data['run_x'], data['tc_y'],  c='green', label='TCase')
            plt.scatter(data['run_x'], data['fb_y'],  c='red',   label='FB')
            plt.title(f'Thermal Profiling – {subfolder}')
            plt.xlabel('run_order')
            plt.ylabel('Temperature')
            plt.legend(loc='upper right', fontsize=8)
            plt.grid(True)
            plt.tight_layout()

            img_path = f'temp_plot_{subfolder}.png'
            plt.savefig(img_path, bbox_inches='tight', dpi=300)
            plt.close()

            # ---- slide ----
            slide = prs.slides.add_slide(slide_layout)
            pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(0.5), width=Inches(9))

            # ---- table ----
            table = slide.shapes.add_table(rows=10, cols=5,
                                          left=Inches(0.5), top=Inches(4.45),
                                          width=Inches(9), height=Inches(3)).table

            # headers
            table.cell(0,0).text = 'Temperature Range'
            table.cell(0,1).text = 'Max_DTS'
            table.cell(0,2).text = 'TCase'
            table.cell(0,3).text = 'FB'
            table.cell(0,4).text = 'Max_DTS %'
            for c in table.rows[0].cells:
                c.text_frame.paragraphs[0].font.size = Pt(10)
                c.text_frame.paragraphs[0].font.bold = True

            dts_stats = calculate_stats(data['dts_y'], temp_ranges)
            tc_stats  = calculate_stats(data['tc_y'],  temp_ranges)
            fb_stats  = calculate_stats(data['fb_y'],  temp_ranges)

            dts_total = sum(dts_stats.values()) or 1   # avoid div-0

            for i, (high, low) in enumerate(temp_ranges, 1):
                r = f"{high}-{low}"
                table.cell(i,0).text = r
                table.cell(i,1).text = str(dts_stats[r])
                table.cell(i,2).text = str(tc_stats[r])
                table.cell(i,3).text = str(fb_stats[r])
                pct = dts_stats[r] / dts_total * 100
                table.cell(i,4).text = f"{pct:.1f}%"
                for c in table.rows[i].cells:
                    c.text_frame.paragraphs[0].font.size = Pt(8)

            os.remove(img_path)

        # ---- SAVE ----
        ts = datetime.now().strftime('%Y%m%d_%H%M%S')
        out_pptx = fr'U:\NVL\HX\A0\results_experimental\ThermalProfileSummary\Thermal_Profiling_Results_{main}_{ts}.pptx'
        os.makedirs(os.path.dirname(out_pptx), exist_ok=True)
        prs.save(out_pptx)
        print(f"\nPowerPoint saved: {out_pptx}")

except Exception as exc:
    print(f"\nFATAL ERROR: {exc}")